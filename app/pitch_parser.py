import os
import fitz  # PyMuPDF
import pptx
from openai import OpenAI

# Load API keys from environment
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

if not OPENROUTER_API_KEY:
    raise ValueError("OPENROUTER_API_KEY environment variable is not set. Please export your OpenRouter API key.")

# Set up OpenRouter client
client = OpenAI(
    api_key=OPENROUTER_API_KEY,
    base_url="https://openrouter.ai/api/v1"
)

# Extract text from PDF
def extract_text_from_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# Extract text from PPTX
def extract_text_from_pptx(path):
    prs = pptx.Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Auto-detect file format and extract text
def extract_text_from_file(file_path):
    """
    Automatically detects file format and extracts text accordingly
    Supports PDF and PPTX formats
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_extension = file_path.lower().split('.')[-1]
    
    if file_extension == 'pdf':
        print(f"[INFO] Detected PDF format, extracting text...")
        return extract_text_from_pdf(file_path)
    elif file_extension in ['pptx', 'ppt']:
        print(f"[INFO] Detected PowerPoint format, extracting text...")
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: PDF, PPTX")


# Extract company information from pitch deck
def extract_company_info(extracted_text):
    print("[INFO] Extracting company information...")
    
    EXTRACTION_PROMPT = (
        "You are an AI assistant that extracts key company information from pitch decks. "
        "Analyze the following pitch deck text and extract ONLY the company name and website URL. "
        "Return your response in this exact JSON format: "
        '{"company_name": "Company Name", "website": "https://website.com"} '
        "If website is not found, use null. Be precise and only extract what's clearly stated."
    )
    
    response = client.chat.completions.create(
        model="openai/gpt-4",
        messages=[
            {"role": "system", "content": EXTRACTION_PROMPT},
            {"role": "user", "content": extracted_text}
        ],
        temperature=0.1
    )
    
    try:
        import json
        result = json.loads(response.choices[0].message.content.strip())
        return result.get("company_name"), result.get("website")
    except:
        # Fallback if JSON parsing fails
        content = response.choices[0].message.content
        return None, None

# Use OpenRouter (ChatGPT or Claude) to generate a markdown summary
def summarize_with_ai(extracted_text, extra_context=""):
    print("[INFO] Sending data to OpenRouter for summarization...")

    SYSTEM_PROMPT = (
        "You are an expert investment analyst for an early-stage VC fund. "
        "Analyze this pitch deck and create a comprehensive investment memo. "
        "Structure your analysis to include:\n"
        "- Company Overview & Mission\n"
        "- Product/Service Description\n"
        "- Market Opportunity & Size\n"
        "- Business Model & Revenue Streams\n"
        "- Traction & Key Metrics\n"
        "- Founding Team & Key Personnel\n"
        "- Competitive Landscape\n"
        "- Financial Projections (if available)\n"
        "- Investment Ask & Use of Funds\n"
        "- Key Risks & Concerns\n"
        "- Investment Recommendation\n\n"
        "Be analytical, concise, and focus on metrics that matter to investors. "
        "Highlight any red flags or areas requiring due diligence."
    )

    full_input = extracted_text + "\n\n" + extra_context

    response = client.chat.completions.create(
        model="openai/gpt-4",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": full_input}
        ],
        temperature=0.3
    )

    return response.choices[0].message.content
